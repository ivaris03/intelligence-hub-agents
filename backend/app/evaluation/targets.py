from __future__ import annotations

import asyncio
import math
import re
from io import BytesIO
from typing import Any, Literal

from langchain_core.messages import HumanMessage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from app.agents.service import (
    _clean_slide_text,
    _normalize_slide_bullets,
    _render_presentation,
)
from app.agents.workflows import (
    DeepResearchHarness,
    PresentationOutline,
    ResearchBudget,
    ResearchResult,
    SlideContent,
    build_research_graph,
    deterministic_outline,
    normalize_presentation_outline,
    research_markdown,
    route_presentation_intent,
)
from app.core.config import Settings
from app.core.security import remove_unverified_urls
from app.files.service import (
    document_chunks_for_upload,
    rerank_document_items,
    validate_upload,
)
from app.integrations.qwen import QwenAdapter
from app.integrations.tavily import (
    TavilyAdapter,
    comprehensive_search,
    normalize_search_citations,
    rerank_search_results,
)

Variant = Literal["baseline", "optimized"]


def _terms(text: str) -> set[str]:
    lower = text.casefold()
    words = set(re.findall(r"[a-z0-9_@.-]{2,}|[\u4e00-\u9fff]", lower))
    chinese = "".join(re.findall(r"[\u4e00-\u9fff]", lower))
    words.update(chinese[index : index + 2] for index in range(max(0, len(chinese) - 1)))
    return words


def _cosine(left: list[float] | None, right: list[float] | None) -> float:
    if not left or not right:
        return 0.0
    numerator = sum(a * b for a, b in zip(left, right, strict=False))
    denominator = math.sqrt(sum(a * a for a in left)) * math.sqrt(sum(b * b for b in right))
    return numerator / denominator if denominator else 0.0


async def rag_target(
    inputs: dict[str, Any],
    settings: Settings,
    variant: Variant,
    attachments: dict[str, Any] | None = None,
) -> dict:
    query = str(inputs["query"])
    documents = list(inputs["documents"])
    uploaded_files: list[dict[str, Any]] = []
    for descriptor in inputs.get("attachment_files") or []:
        key = str(descriptor["key"])
        attachment = (attachments or {}).get(key)
        if attachment is None:
            raise ValueError(f"评估样本缺少附件：{key}")
        reader = attachment.get("reader")
        if reader is None:
            raise ValueError(f"评估附件没有可读内容：{key}")
        data = reader.read()
        validated = validate_upload(
            str(descriptor["filename"]),
            str(attachment.get("mime_type") or "application/octet-stream"),
            data,
            settings,
        )
        chunks = document_chunks_for_upload(validated, settings)
        base_id = str(descriptor["document_id"])
        evaluation_chunks = chunks[:32]
        documents.extend(
            {
                "id": base_id if index == 0 else f"{base_id}:{index}",
                "text": chunk.content,
            }
            for index, chunk in enumerate(evaluation_chunks)
        )
        uploaded_files.append(
            {
                "name": descriptor["filename"],
                "size_bytes": len(data),
                "kind": validated.kind,
                "chunks": len(chunks),
            }
        )
    texts = [query, *(str(item["text"]) for item in documents)]
    try:
        embeddings = await QwenAdapter(settings).embed_documents(texts)
    except Exception:
        embeddings = [None] * len(texts)
    query_embedding = embeddings[0]
    query_terms = _terms(query)
    scored: list[tuple[float, dict[str, str]]] = []
    for item, embedding in zip(documents, embeddings[1:], strict=True):
        semantic = _cosine(query_embedding, embedding)
        lexical = len(query_terms & _terms(item["text"])) / max(1, len(query_terms))
        score = semantic if variant == "baseline" else semantic * 0.65 + lexical * 0.35
        scored.append((score, item))
    ranked = [item for _, item in sorted(scored, key=lambda pair: pair[0], reverse=True)]
    if variant == "optimized":
        candidates = await rerank_document_items(
            query,
            [{"text": item["text"], "document": item} for item in ranked[:8]],
            settings,
            limit=5,
        )
        retrieved = [item["document"] for item in candidates]
    else:
        retrieved = ranked[:5]
    context = "\n\n".join(f"[{item['id']}] {item['text']}" for item in retrieved)
    if variant == "baseline":
        instruction = "根据下面检索资料回答问题。"
    else:
        instruction = (
            "只能根据检索资料直接回答用户问到的内容，不要补充问题未问到的尾注或"
            "‘其他信息资料未提供’；确实无法回答时才写‘资料未提供’。"
            "不得把只适用于部分实体的规则泛化为全部实体。"
            "答案尽量简短，每个事实结论后标注对应的 [文档ID]。"
        )
    generation_mode = "qwen"
    try:
        answer = await QwenAdapter(settings).complete_text(
            f"{instruction}\n\n检索资料：\n{context}\n\n问题：{query}",
            fallback="资料未提供。",
            max_chars=3_000,
        )
    except Exception:
        generation_mode = "extractive_fallback"
        primary = retrieved[0] if retrieved else None
        answer = (
            f"{primary['text']} [{primary['id']}]" if primary else "资料未提供。"
        )
    return {
        "answer": answer,
        "retrieved_document_ids": [item["id"] for item in retrieved],
        "retrieved_documents": retrieved,
        "uploaded_files": uploaded_files,
        "generation_mode": generation_mode,
        "variant": variant,
    }


def clean_search_query(text: str) -> str:
    cleaned = re.sub(
        r"^(?:请|帮我|麻烦)?\s*(?:联网|上网|网上)?\s*(?:搜索|查找|查询|查|搜)\s*(?:一下)?\s*",
        "",
        text.strip(),
        flags=re.I,
    )
    cleaned = re.sub(
        r"(?i)^(?:please\s+)?(?:search|browse|look up)\s+"
        r"(?:the\s+)?(?:web|internet|online)\s*(?:for)?\s*",
        "",
        cleaned,
    )
    return cleaned.strip(" ：:，,") or text.strip()


def _extractive_search_answer(results: list[Any]) -> str:
    """Build a citation-safe answer when generation is temporarily unavailable."""

    lines = ["根据本次可核验搜索结果："]
    for item in results:
        snippet = re.sub(r"\s+", " ", str(item.snippet or "")).strip()
        if not snippet:
            continue
        sentence = re.split(r"(?<=[。！？.!?])\s+", snippet, maxsplit=1)[0]
        sentence = sentence[:280].rstrip("，,；;：: ")
        lines.append(f"- {sentence} ([{item.title}]({item.url}))")
    return "\n".join(lines) if len(lines) > 1 else "搜索结果没有返回可引用的摘要。"


async def mcp_target(inputs: dict[str, Any], settings: Settings, variant: Variant) -> dict:
    user_query = str(inputs["query"])
    search_query = user_query if variant == "baseline" else clean_search_query(user_query)
    candidates = (
        await TavilyAdapter(settings).search(search_query, max_results=5)
        if variant == "baseline"
        else await comprehensive_search(search_query, settings)
    )
    results = (
        candidates
        if variant == "baseline"
        else await rerank_search_results(search_query, candidates, settings, limit=5)
    )
    source_text = "\n".join(
        f"[{index}] {item.title}\nURL: {item.url}\n摘要: {item.snippet}"
        for index, item in enumerate(results, 1)
    )
    if variant == "baseline":
        instruction = "根据搜索结果回答，并且只使用搜索结果中的 URL。"
    else:
        instruction = (
            "只根据搜索结果回答。先覆盖问题的各个方面，再精炼结论；"
            "每个可核验事实后都用 Markdown 链接引用直接支持它的来源；"
            "链接必须逐字使用搜索结果中的 URL，摘要不支持的结论不要写。"
        )
    generation_mode = "qwen"
    try:
        answer = await QwenAdapter(settings).complete_text(
            f"{instruction}\n\n用户问题：{user_query}\n\n搜索结果：\n{source_text}",
            fallback="搜索结果不足，暂时无法回答。",
            work=True,
            max_chars=5_000,
        )
    except Exception:
        generation_mode = "extractive_fallback"
        answer = _extractive_search_answer(results)
    if variant == "optimized":
        answer = normalize_search_citations(answer, results)
        answer, _ = remove_unverified_urls(answer, {item.url for item in results})
        if results and not any(item.url in answer for item in results):
            primary = results[0]
            answer += f"\n\n来源：[{primary.title}]({primary.url})"
    return {
        "answer": answer,
        "search_query": search_query,
        "search_results": [item.as_dict() for item in results],
        "generation_mode": generation_mode,
        "variant": variant,
    }


def _legacy_route_presentation_intent(
    text: str,
    *,
    requested: str | None = None,
    source_artifact_id: str | None = None,
    source_run_id: str | None = None,
) -> str:
    if requested in {"CREATE", "MODIFY", "RESUME"}:
        intent = requested
    elif source_run_id or re.search(
        r"^(?:请\s*)?(?:(?:继续|恢复)(?:$|这|该|上|原|之前|运行|任务|演示|生成)|resume\b)"
        r"|从.+(?:继续|恢复)",
        text.strip(),
        re.I,
    ):
        intent = "RESUME"
    elif source_artifact_id or re.search(r"修改|调整|替换|改成|modify|edit", text, re.I):
        intent = "MODIFY"
    else:
        intent = "CREATE"
    if intent == "MODIFY" and not source_artifact_id:
        raise ValueError("MODIFY 必须指定源演示版本")
    if intent == "RESUME" and not source_run_id:
        raise ValueError("RESUME 必须指定原运行")
    return intent


def _legacy_render_presentation(outline: PresentationOutline, slides: list[SlideContent]) -> bytes:
    presentation = Presentation()
    presentation.slide_width = 12_191_695
    presentation.slide_height = 6_858_000
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(247, 245, 239)
    title_slide.shapes.title.text = outline.title
    title_slide.placeholders[1].text = "Intelligence Hub · AI 生成草稿"
    for item in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(251, 250, 247)
        slide.shapes.title.text = item.title
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.color.rgb = RGBColor(39, 95, 84)
        title_frame.paragraphs[0].font.size = Pt(28)
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(item.bullets[:6]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(20)
            paragraph.space_after = Pt(12)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _rgb(value: Any) -> str | None:
    try:
        rgb = value.rgb
        return str(rgb) if rgb is not None else None
    except (AttributeError, TypeError, ValueError):
        return None


def presentation_manifest(data: bytes) -> dict[str, Any]:
    presentation = Presentation(BytesIO(data))
    slides: list[dict[str, Any]] = []
    for number, slide in enumerate(presentation.slides, 1):
        shapes = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if not text:
                continue
            frame = getattr(shape, "text_frame", None)
            paragraphs = list(frame.paragraphs) if frame is not None else []
            font_sizes = [
                round(size.pt, 1)
                for paragraph in paragraphs
                for size in [
                    paragraph.font.size,
                    *(run.font.size for run in paragraph.runs),
                ]
                if size is not None
            ]
            font_colors = [
                color
                for paragraph in paragraphs
                for font in [paragraph.font, *(run.font for run in paragraph.runs)]
                if (color := _rgb(font.color))
            ]
            shapes.append(
                {
                    "text": text[:1_200],
                    "x": round(shape.left / presentation.slide_width, 3),
                    "y": round(shape.top / presentation.slide_height, 3),
                    "w": round(shape.width / presentation.slide_width, 3),
                    "h": round(shape.height / presentation.slide_height, 3),
                    "font_sizes_pt": font_sizes,
                    "font_colors": font_colors,
                }
            )
        slides.append(
            {
                "number": number,
                "background": _rgb(slide.background.fill.fore_color),
                "shapes": shapes,
            }
        )
    return {
        "slide_count": len(slides),
        "aspect_ratio": round(presentation.slide_width / presentation.slide_height, 3),
        "slides": slides,
    }


async def _generate_presentation(
    inputs: dict[str, Any], settings: Settings, variant: Variant
) -> dict[str, Any]:
    task = str(inputs["text"])
    audience = str(inputs.get("audience") or "一般读者")
    max_total_pages = int(inputs.get("max_total_pages") or 6)
    max_content_pages = max(2, max_total_pages - 1)
    fallback = deterministic_outline(task, max_content_pages)
    model = QwenAdapter(settings).chat_model(work=True)
    outline_model = model.with_structured_output(PresentationOutline)
    if variant == "baseline":
        outline_instruction = (
            "为下面的演示任务生成结构化大纲。页面标题应互不重复，不得使用占位文本。"
        )
    else:
        outline_instruction = (
            "你是资深演示叙事设计师。为任务生成结构化大纲：先建立背景和决策问题，"
            "再给证据/分析，最后给明确行动；每一页只承担一个角色，标题要表达结论而不是主题词；"
            "页面互不重复，并完整覆盖用户的显式要求。slides 中每个元素只能是纯文本标题，"
            "不得把 JSON 对象序列化成字符串，也不得编造用户未提供的数字或业务结果。"
        )
    try:
        generated_outline = await outline_model.ainvoke(
            [
                HumanMessage(
                    content=(
                        f"{outline_instruction}\n受众：{audience}\n"
                        f"含封面不超过 {max_total_pages} 页。\n任务：{task}"
                    )
                )
            ]
        )
        outline = normalize_presentation_outline(
            PresentationOutline.model_validate(generated_outline), task, max_content_pages
        )
    except Exception:
        outline = fallback

    content_model = model.with_structured_output(SlideContent)
    slides: list[SlideContent] = []
    for index, title in enumerate(outline.slides):
        if variant == "baseline":
            prompt = (
                f"为演示任务“{task}”生成这一页的结构化内容。页标题必须保持为“{title}”。"
                "要点不超过 5 条，不得编造未提供的数据；缺失信息写‘待补充’。"
            )
        else:
            prompt = (
                "你是资深演示编辑。生成可直接上屏而非讲稿式的页面内容：每条要点简短、具体、"
                "彼此不重复，优先使用‘结论—依据—行动’关系；不得编造数据。"
                f"\n任务：{task}\n受众：{audience}\n完整大纲：{outline.slides}"
                f"\n当前为第 {index + 1}/{len(outline.slides)} 页，标题必须保持为“{title}”。"
                "\n未提供的精确值改写为验证方法或衡量计划，不要堆砌‘待补充’。"
                "给 3–4 条纯文本要点，每条不超过 90 字，并在 speaker_notes 中说明本页目的和转场。"
            )
        fallback_slide = SlideContent(
            title=title,
            bullets=[f"围绕“{outline.title}”说明 {title}", "关键事实与建议", "可执行的下一步"],
        )
        try:
            generated = SlideContent.model_validate(
                await content_model.ainvoke([HumanMessage(content=prompt)])
            )
            bullets = _normalize_slide_bullets(
                title, generated.bullets, fallback_slide.bullets
            )
            slides.append(
                SlideContent(
                    title=title,
                    bullets=bullets,
                    speaker_notes=_clean_slide_text(generated.speaker_notes),
                )
            )
        except Exception:
            slides.append(
                SlideContent(
                    title=title,
                    bullets=_normalize_slide_bullets(title, fallback_slide.bullets),
                )
            )

    if variant == "baseline":
        pptx_data = _legacy_render_presentation(outline, slides)
    else:
        pptx_data, _ = _render_presentation(outline, slides)
    return {
        "outline": outline.model_dump(),
        "slides": [slide.model_dump() for slide in slides],
        "design_manifest": presentation_manifest(pptx_data),
        "artifact_size_bytes": len(pptx_data),
    }


async def slides_target(inputs: dict[str, Any], settings: Settings, variant: Variant) -> dict:
    if inputs["task"] == "intent":
        router = (
            _legacy_route_presentation_intent
            if variant == "baseline"
            else route_presentation_intent
        )
        try:
            intent = router(
                str(inputs["text"]),
                requested=inputs.get("requested"),
                source_artifact_id="artifact" if inputs.get("has_source_artifact") else None,
                source_run_id="run" if inputs.get("has_source_run") else None,
            )
        except ValueError:
            intent = "ERROR"
        return {"task": "intent", "intent": intent, "variant": variant}
    generated = await _generate_presentation(inputs, settings, variant)
    return {"task": "quality", "variant": variant, **generated}


async def research_target(inputs: dict[str, Any], settings: Settings, variant: Variant) -> dict:
    budget = ResearchBudget(
        max_searches=settings.research_max_searches,
        timeout_seconds=settings.research_timeout_seconds,
    )
    harness = DeepResearchHarness(settings, budget, prompt_version=variant)
    graph = build_research_graph(harness)
    outer_timeout = (
        settings.research_timeout_seconds
        if variant == "baseline"
        else settings.research_timeout_seconds + 3
    )
    async with asyncio.timeout(outer_timeout):
        state = await graph.ainvoke(
            {"question": str(inputs["question"]), "context": ""},
            config={
                "run_name": "evaluation.deep_research",
                "tags": ["evaluation", f"variant:{variant}"],
            },
        )
    result = ResearchResult.model_validate(state["result"])
    return {
        "report": research_markdown(result),
        "result": result.model_dump(mode="json"),
        "searches_used": budget.used_searches,
        "generation_mode": harness.synthesis_mode,
        "validation_errors": state.get("validation_errors") or [],
        "variant": variant,
    }


def target_for(suite: str, settings: Settings, variant: Variant):
    targets = {
        "rag": rag_target,
        "mcp": mcp_target,
        "slides": slides_target,
        "research": research_target,
    }
    target = targets[suite]

    if suite == "rag":

        async def invoke(inputs: dict[str, Any], attachments: dict[str, Any]) -> dict[str, Any]:
            return await rag_target(inputs, settings, variant, attachments)

    else:

        async def invoke(inputs: dict[str, Any]) -> dict[str, Any]:
            return await target(inputs, settings, variant)

    invoke.__name__ = f"{suite}_{variant}_target"
    return invoke
