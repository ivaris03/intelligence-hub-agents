import asyncio
import time
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine
from sqlalchemy.pool import NullPool

from app.auth.security import get_current_user, hash_password
from app.core.config import Settings, get_settings
from app.db.base import Base, User
from app.db.session import get_session
from app.integrations.qwen import QwenAdapter
from app.main import app


@pytest.fixture
def mvp_client():
    app.dependency_overrides.pop(get_current_user, None)
    temporary = TemporaryDirectory(prefix=".mvp-test-", dir=Path.cwd())
    tmp_path = Path(temporary.name)
    database_path = tmp_path / "mvp.sqlite3"
    database_url = f"sqlite+aiosqlite:///{database_path.as_posix()}"
    engine = create_async_engine(database_url, poolclass=NullPool)
    sessions = async_sessionmaker(engine, expire_on_commit=False)
    settings = Settings(
        database_url=database_url,
        storage_backend="local",
        storage_path=tmp_path / "storage",
        dashscope_api_key=None,
        tavily_api_key=None,
        langsmith_tracing=False,
        langsmith_api_key=None,
    )

    async def create_schema() -> None:
        async with engine.begin() as connection:
            await connection.run_sync(Base.metadata.create_all)
        async with sessions() as session:
            encoded_password = hash_password(
                "12345678", salt=b"integration-test"
            )
            session.add_all(
                [
                User(
                    phone="13900000001",
                    password_hash=encoded_password,
                    display_name="测试管理员",
                    role="admin",
                    is_active=True,
                ),
                User(
                    phone="13700000001",
                    password_hash=encoded_password,
                    display_name="测试用户",
                    role="member",
                    is_active=True,
                ),
                ]
            )
            await session.commit()

    async def override_session():
        async with sessions() as session:
            yield session

    asyncio.run(create_schema())
    app.dependency_overrides[get_session] = override_session
    app.dependency_overrides[get_settings] = lambda: settings
    client = TestClient(app)
    login = client.post(
        "/api/auth/login",
        json={"phone": "13900000001", "password": "12345678"},
    )
    assert login.status_code == 200
    client.headers["Authorization"] = f"Bearer {login.json()['access_token']}"
    try:
        yield client
    finally:
        client.close()
        app.dependency_overrides.clear()
        asyncio.run(engine.dispose())
        temporary.cleanup()


def create_conversation(client: TestClient, mode: str = "chat") -> str:
    response = client.post("/api/conversations", json={"mode": mode})
    assert response.status_code == 201
    return response.json()["id"]


def test_auth_rbac_and_workspace_isolation(mvp_client: TestClient) -> None:
    admin_authorization = mvp_client.headers["Authorization"]
    admin_conversation_id = create_conversation(mvp_client)

    unauthenticated = mvp_client.get(
        "/api/conversations", headers={"Authorization": ""}
    )
    assert unauthenticated.status_code == 401

    login = mvp_client.post(
        "/api/auth/login",
        json={"phone": "13700000001", "password": "12345678"},
    )
    assert login.status_code == 200
    member_authorization = f"Bearer {login.json()['access_token']}"
    member_headers = {"Authorization": member_authorization}
    member_conversation = mvp_client.post(
        "/api/conversations", json={"mode": "chat"}, headers=member_headers
    )
    assert member_conversation.status_code == 201

    member_list = mvp_client.get("/api/conversations", headers=member_headers)
    assert [item["id"] for item in member_list.json()] == [
        member_conversation.json()["id"]
    ]
    assert mvp_client.get(
        f"/api/conversations/{admin_conversation_id}/messages", headers=member_headers
    ).status_code == 404
    assert login.json()["user"]["role"] == "member"

    admin_headers = {"Authorization": admin_authorization}
    assert mvp_client.get(
        f"/api/conversations/{member_conversation.json()['id']}/messages",
        headers=admin_headers,
    ).status_code == 404
    administrator = mvp_client.get("/api/auth/me", headers=admin_headers)
    assert administrator.status_code == 200
    assert administrator.json()["phone"] == "13900000001"
    assert administrator.json()["role"] == "admin"


def test_single_memory_summary_api_and_commands(mvp_client: TestClient) -> None:
    initial = mvp_client.get("/api/memory-summary")
    assert initial.status_code == 200
    assert initial.json()["id"] == 1
    assert initial.json()["content"] == ""

    updated = mvp_client.put(
        "/api/memory-summary", json={"content": "用户正在开发 Intelligence Hub。"}
    )
    assert updated.status_code == 200
    assert updated.json()["content"] == "用户正在开发 Intelligence Hub。"

    conversation_id = create_conversation(mvp_client)
    remembered = mvp_client.post(
        f"/api/conversations/{conversation_id}/messages",
        json={"content": "请记住我常用 Python", "mode": "chat"},
    )
    assert "event: memory.updated" in remembered.text
    summary = mvp_client.get("/api/memory-summary").json()
    assert summary["id"] == 1
    assert summary["content"] == "用户正在开发 Intelligence Hub；我常用 Python。"

    cleared = mvp_client.delete("/api/memory-summary")
    assert cleared.status_code == 204
    assert mvp_client.get("/api/memory-summary").json()["content"] == ""


def test_memory_summary_chat_can_answer_and_replace_memory(mvp_client: TestClient) -> None:
    updated = mvp_client.put(
        "/api/memory-summary", json={"content": "用户喜欢苹果。"}
    )
    assert updated.status_code == 200

    question = mvp_client.post(
        "/api/memory-summary/messages", json={"content": "你记住了什么？"}
    )
    assert question.status_code == 200
    assert question.json()["changed"] is False
    assert "用户喜欢苹果" in question.json()["assistant_message"]["content"]

    correction = mvp_client.post(
        "/api/memory-summary/messages", json={"content": "我现在喜欢吃梨了"}
    )
    assert correction.status_code == 200
    assert correction.json()["changed"] is True
    assert correction.json()["summary"]["content"] == "我喜欢吃梨。"
    assert correction.json()["assistant_message"]["memory_changed"] is True

    history = mvp_client.get("/api/memory-summary/messages")
    assert history.status_code == 200
    assert [item["role"] for item in history.json()] == [
        "user",
        "assistant",
        "user",
        "assistant",
    ]

    rejected = mvp_client.post(
        "/api/memory-summary/messages",
        json={"content": "我的 API_KEY 是 sk-example123456"},
    )
    assert rejected.status_code == 200
    assert rejected.json()["changed"] is False
    assert "不会" in rejected.json()["assistant_message"]["content"]
    assert rejected.json()["summary"]["content"] == "我喜欢吃梨。"


def test_manual_memory_refinement_uses_unprocessed_messages(mvp_client: TestClient) -> None:
    conversation_id = create_conversation(mvp_client)
    response = mvp_client.post(
        f"/api/conversations/{conversation_id}/messages",
        json={"content": "我常用 Rust", "mode": "chat"},
    )
    assert response.status_code == 200

    refined = mvp_client.post("/api/memory-summary/refine")
    assert refined.status_code == 200
    assert refined.json()["processed_messages"] == 1
    assert refined.json()["added_facts"] == 1
    assert refined.json()["summary"]["content"] == "我常用 Rust。"

    repeated = mvp_client.post("/api/memory-summary/refine")
    assert repeated.status_code == 200
    assert repeated.json()["processed_messages"] == 0
    assert repeated.json()["added_facts"] == 0


def test_persistent_chat_files_skill_memory_and_regeneration(mvp_client: TestClient) -> None:
    conversation_id = create_conversation(mvp_client)
    renamed = mvp_client.patch(
        f"/api/conversations/{conversation_id}", json={"title": "手动标题"}
    )
    assert renamed.status_code == 200
    assert renamed.json()["title_source"] == "manual"
    skill = mvp_client.post(
        "/api/skills",
        json={
            "name": "简洁写作",
            "description": "简短概括材料",
            "instructions": "先给结论，再给三个要点。",
            "enabled": True,
        },
    ).json()
    uploaded = mvp_client.post(
        "/api/files",
        data={"conversation_id": conversation_id},
        files={"upload": ("notes.txt", "项目目标是完成可演示的 Chat 和 Agent MVP。", "text/plain")},
    )
    assert uploaded.status_code == 201
    file_id = uploaded.json()["id"]

    memory_stream = mvp_client.post(
        f"/api/conversations/{conversation_id}/messages",
        json={"content": "你好，请记住我偏好简洁回答", "mode": "chat"},
    )
    assert "event: memory.updated" in memory_stream.text
    assert mvp_client.get("/api/memory-summary").json()["content"] == "我偏好简洁回答。"

    response = mvp_client.post(
        f"/api/conversations/{conversation_id}/messages",
        json={
            "content": "概括项目目标",
            "mode": "chat",
            "file_ids": [file_id],
            "skill_id": skill["id"],
        },
    )
    assert response.status_code == 200
    assert "event: sources.finalized" in response.text
    assert "event: skill.selected" in response.text
    assert "event: follow_up.finalized" in response.text

    messages = mvp_client.get(f"/api/conversations/{conversation_id}/messages").json()
    assert [message["role"] for message in messages[:2]] == ["user", "assistant"]
    assistant = messages[-1]
    assert assistant["files"][0]["id"] == file_id
    assert assistant["skill"]["name"] == "简洁写作"
    assert assistant["follow_up"]
    user_count = sum(message["role"] == "user" for message in messages)

    regenerated = mvp_client.post(f"/api/messages/{assistant['id']}/regenerate")
    assert regenerated.status_code == 200
    refreshed = mvp_client.get(f"/api/conversations/{conversation_id}/messages").json()
    assert sum(message["role"] == "user" for message in refreshed) == user_count
    assert refreshed[-1]["files"][0]["id"] == file_id

    search = mvp_client.get("/api/conversations", params={"q": "项目目标"}).json()
    assert search and search[0]["match_snippet"]


def test_skill_mentions_are_ignored_and_work_accepts_multiple_skills(
    mvp_client: TestClient,
) -> None:
    skill_ids: list[str] = []
    for name, instruction in (
        ("简洁写作", "先给结论。"),
        ("风险检查", "补充主要风险。"),
    ):
        response = mvp_client.post(
            "/api/skills",
            json={
                "name": name,
                "description": instruction,
                "instructions": instruction,
                "enabled": True,
            },
        )
        assert response.status_code == 201
        skill_ids.append(response.json()["id"])

    chat_id = create_conversation(mvp_client)
    mentioned = mvp_client.post(
        f"/api/conversations/{chat_id}/messages",
        json={"content": "@简洁写作 这只是普通文字", "mode": "chat"},
    )
    assert mentioned.status_code == 200
    assert "event: skill.selected" not in mentioned.text
    chat_messages = mvp_client.get(f"/api/conversations/{chat_id}/messages").json()
    assert chat_messages[-1]["skill"] is None

    automatic = mvp_client.post(
        f"/api/conversations/{chat_id}/messages",
        json={"content": "请做风险检查", "mode": "chat"},
    )
    assert automatic.status_code == 200
    assert "event: skill.selected" in automatic.text
    automatic_message = mvp_client.get(
        f"/api/conversations/{chat_id}/messages"
    ).json()[-1]
    assert [skill["name"] for skill in automatic_message["skills"]] == ["风险检查"]

    selected = mvp_client.post(
        f"/api/conversations/{chat_id}/messages",
        json={
            "content": "同时应用两个 Skill",
            "mode": "chat",
            "skill_ids": skill_ids,
        },
    )
    assert selected.status_code == 200
    assert "event: skills.selected" in selected.text
    chat_messages = mvp_client.get(f"/api/conversations/{chat_id}/messages").json()
    assert [skill["name"] for skill in chat_messages[-1]["skills"]] == [
        "简洁写作",
        "风险检查",
    ]
    regenerated = mvp_client.post(f"/api/messages/{chat_messages[-1]['id']}/regenerate")
    assert regenerated.status_code == 200
    regenerated_message = mvp_client.get(
        f"/api/conversations/{chat_id}/messages"
    ).json()[-1]
    assert [skill["name"] for skill in regenerated_message["skills"]] == [
        "简洁写作",
        "风险检查",
    ]

    work_id = create_conversation(mvp_client, "work")
    started = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": work_id,
            "agent_type": "image",
            "input": "生成一张知识中心插画",
            "skill_ids": skill_ids,
        },
    )
    assert started.status_code == 200
    run = mvp_client.get(f"/api/conversations/{work_id}/agent-runs").json()[0]
    assert [skill["name"] for skill in run["skills"]] == ["简洁写作", "风险检查"]

    multiple_agents = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": work_id,
            "agent_type": ["image", "slides"],
            "input": "不能同时选择两个 Agent",
        },
    )
    assert multiple_agents.status_code == 422


def test_chat_and_work_accept_per_request_thinking_effort(
    mvp_client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    observed_efforts: list[str] = []

    async def record_stream(self, messages, **kwargs):
        del messages, kwargs
        observed_efforts.append(self.settings.qwen_thinking_effort)
        yield "text", "完成"

    monkeypatch.setattr(QwenAdapter, "stream_chat", record_stream)
    conversation_id = create_conversation(mvp_client)
    response = mvp_client.post(
        f"/api/conversations/{conversation_id}/messages",
        json={"content": "深入分析", "mode": "chat", "thinking_effort": "high"},
    )
    assert response.status_code == 200
    assert observed_efforts == ["high"]

    assistant = mvp_client.get(
        f"/api/conversations/{conversation_id}/messages"
    ).json()[-1]
    regenerated = mvp_client.post(
        f"/api/messages/{assistant['id']}/regenerate",
        json={"thinking_effort": "none"},
    )
    assert regenerated.status_code == 200
    assert observed_efforts[-1] == "none"

    work_id = create_conversation(mvp_client, "work")
    work = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": work_id,
            "agent_type": "image",
            "input": "生成一张图",
            "thinking_effort": "high",
        },
    )
    assert work.status_code == 200
    runs = mvp_client.get(f"/api/conversations/{work_id}/agent-runs").json()
    assert runs[0]["public_state"]["thinking_effort"] == "high"


def test_chat_and_work_use_separate_conversations(mvp_client: TestClient) -> None:
    chat_id = create_conversation(mvp_client, "chat")
    work_id = create_conversation(mvp_client, "work")

    chat_conversations = mvp_client.get("/api/conversations", params={"mode": "chat"}).json()
    work_conversations = mvp_client.get("/api/conversations", params={"mode": "work"}).json()
    all_conversations = mvp_client.get("/api/conversations").json()
    assert {item["id"] for item in chat_conversations} == {chat_id}
    assert {item["id"] for item in work_conversations} == {work_id}
    assert {item["id"] for item in all_conversations} == {chat_id, work_id}
    assert all(item["mode"] == "chat" for item in chat_conversations)
    assert all(item["mode"] == "work" for item in work_conversations)

    chat_in_work = mvp_client.post(
        f"/api/conversations/{work_id}/messages",
        json={"content": "这条消息不能进入 Work", "mode": "chat"},
    )
    assert chat_in_work.status_code == 409
    assert chat_in_work.json()["detail"] == "Chat 和 Work 不能共用同一个会话"

    work_in_chat = mvp_client.post(
        "/api/agent-runs",
        json={"conversation_id": chat_id, "agent_type": "image", "input": "不能运行"},
    )
    assert work_in_chat.status_code == 400
    assert work_in_chat.json()["detail"] == "Work 任务只能写入 Work 会话"


def test_all_agent_artifacts_and_slide_confirmation(mvp_client: TestClient) -> None:
    conversation_id = create_conversation(mvp_client, "work")
    image_stream = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": conversation_id,
            "agent_type": "image",
            "input": "绿色知识中心插画",
        },
    )
    assert "event: artifact.created" in image_stream.text

    slide_stream = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": conversation_id,
            "agent_type": "slides",
            "input": "制作项目演示",
        },
    )
    assert "event: outline.ready" in slide_stream.text
    assert "awaiting_confirmation" in slide_stream.text
    runs = mvp_client.get(f"/api/conversations/{conversation_id}/agent-runs").json()
    slide_run = next(run for run in runs if run["agent_type"] == "slides")
    assert slide_run["artifacts"] == []
    premature_resume = mvp_client.post(f"/api/agent-runs/{slide_run['id']}/resume")
    assert premature_resume.status_code == 409
    confirmed = mvp_client.post(
        f"/api/agent-runs/{slide_run['id']}/commands", json={"action": "confirm"}
    )
    assert "event: artifact.created" in confirmed.text
    final_slide_run = mvp_client.get(f"/api/agent-runs/{slide_run['id']}").json()
    pptx = final_slide_run["artifacts"][0]
    download = mvp_client.get(pptx["download_url"])
    assert download.content.startswith(b"PK")

    modified_start = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": conversation_id,
            "agent_type": "slides",
            "input": "只修改第 2 页，补充验收说明",
            "intent": "MODIFY",
            "source_artifact_id": pptx["id"],
        },
    )
    assert "awaiting_confirmation" in modified_start.text
    modified_run = mvp_client.get(f"/api/conversations/{conversation_id}/agent-runs").json()[-1]
    assert modified_run["public_state"]["modification_plan"]["target_slides"] == [2]
    mvp_client.post(
        f"/api/agent-runs/{modified_run['id']}/commands", json={"action": "confirm"}
    )
    modified_run = mvp_client.get(f"/api/agent-runs/{modified_run['id']}").json()
    modified = modified_run["artifacts"][0]
    assert modified["version"] == 2
    assert modified["parent_artifact_id"] == pptx["id"]
    assert mvp_client.get(pptx["download_url"]).content == download.content
    modified_deck = Presentation(BytesIO(mvp_client.get(modified["download_url"]).content))
    assert "已更新" in modified_deck.slides[1].shapes.title.text
    assert "已更新" not in modified_deck.slides[2].shapes.title.text

    interrupted = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": conversation_id,
            "agent_type": "slides",
            "input": "制作可恢复演示",
        },
    )
    assert "awaiting_confirmation" in interrupted.text
    interrupted_run = mvp_client.get(
        f"/api/conversations/{conversation_id}/agent-runs"
    ).json()[-1]
    cancelled = mvp_client.post(
        f"/api/agent-runs/{interrupted_run['id']}/commands", json={"action": "cancel"}
    )
    assert cancelled.json()["status"] == "cancelled"
    resumed = mvp_client.post(f"/api/agent-runs/{interrupted_run['id']}/resume")
    assert "event: artifact.created" in resumed.text
    resumed_run = mvp_client.get(f"/api/agent-runs/{interrupted_run['id']}").json()
    assert resumed_run["status"] == "completed"
    assert len(resumed_run["artifacts"]) == 1
    mvp_client.post(f"/api/agent-runs/{interrupted_run['id']}/resume")
    assert len(mvp_client.get(f"/api/agent-runs/{interrupted_run['id']}").json()["artifacts"]) == 1

    research = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": conversation_id,
            "agent_type": "research",
            "input": "研究个人 Agent Hub 的价值",
        },
    )
    assert "event: artifact.created" in research.text
    runs = mvp_client.get(f"/api/conversations/{conversation_id}/agent-runs").json()
    assert {artifact["type"] for run in runs for artifact in run["artifacts"]} == {
        "image",
        "pptx",
        "markdown",
    }


def test_agent_failure_is_visible_and_retryable(mvp_client: TestClient, monkeypatch) -> None:
    conversation_id = create_conversation(mvp_client, "work")
    attempts = 0

    async def flaky_generate(self, prompt, reference_images=None):
        nonlocal attempts
        attempts += 1
        if attempts == 1:
            raise RuntimeError("provider payload must not reach the client")
        return self._demo_image()

    monkeypatch.setattr(QwenAdapter, "generate_image", flaky_generate)
    failed = mvp_client.post(
        "/api/agent-runs",
        json={
            "conversation_id": conversation_id,
            "agent_type": "image",
            "input": "测试可恢复失败",
        },
    )
    assert "event: tool.failed" in failed.text
    assert "provider payload" not in failed.text
    run = mvp_client.get(f"/api/conversations/{conversation_id}/agent-runs").json()[0]
    assert (run["status"], run["stage"]) == ("failed", "failed")
    retried = mvp_client.post(
        f"/api/agent-runs/{run['id']}/commands", json={"action": "retry"}
    )
    assert "event: artifact.created" in retried.text
    final = mvp_client.get(f"/api/agent-runs/{run['id']}").json()
    assert final["status"] == "completed"
    assert len(final["artifacts"]) == 1


def test_stream_can_stop_without_duplicating_user_message(
    mvp_client: TestClient, monkeypatch
) -> None:
    conversation_id = create_conversation(mvp_client)

    async def slow_stream(self, messages, **kwargs):
        del self, messages, kwargs
        for _ in range(200):
            await asyncio.sleep(0.005)
            yield "text", "片"

    monkeypatch.setattr(QwenAdapter, "stream_chat", slow_stream)
    with ThreadPoolExecutor(max_workers=1) as pool:
        request = pool.submit(
            mvp_client.post,
            f"/api/conversations/{conversation_id}/messages",
            json={"content": "测试停止生成", "mode": "chat"},
        )
        assistant_id = None
        for _ in range(80):
            messages = mvp_client.get(
                f"/api/conversations/{conversation_id}/messages"
            ).json()
            streaming = [
                message
                for message in messages
                if message["role"] == "assistant" and message["status"] == "streaming"
            ]
            if streaming:
                assistant_id = streaming[0]["id"]
                break
            time.sleep(0.01)
        assert assistant_id is not None
        stopped = mvp_client.post(f"/api/messages/{assistant_id}/stop")
        assert stopped.status_code == 200
        assert stopped.json()["status"] == "cancelled"
        stream = request.result(timeout=5)
    assert "event: cancelled" in stream.text
    messages = mvp_client.get(f"/api/conversations/{conversation_id}/messages").json()
    assert sum(message["role"] == "user" for message in messages) == 1
    assert messages[-1]["status"] == "cancelled"
